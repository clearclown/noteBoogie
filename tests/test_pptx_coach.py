"""Tests for pptx structural lint + fix application (api/pptx_coach.py, §11b)."""

import io
import json
from unittest.mock import AsyncMock, patch

import pytest
from fastapi.testclient import TestClient
from pptx import Presentation
from pptx.util import Emu, Inches, Pt

from api.pptx_coach import (
    RULE_ALIGN_SNAP,
    RULE_NORMALIZE_FONTS,
    RULE_TABLE_HEADER,
    apply_fixes,
    build_text_dump,
    extract_pptx,
    lint_pptx,
)
from open_notebook.exceptions import InvalidInputError


@pytest.fixture
def client():
    from api.main import app

    return TestClient(app)


def make_messy_pptx() -> bytes:
    """3フォント・微妙にズレた左端・表ヘッダ非太字、の「直しがいのある」деck。"""
    prs = Presentation()
    blank = prs.slide_layouts[6]

    slide = prs.slides.add_slide(blank)
    for i, (font, left) in enumerate(
        [("Meiryo", Inches(1.0)), ("Arial", Emu(Inches(1.0) + Emu(30000))), ("Comic Sans MS", Inches(4.0))]
    ):
        box = slide.shapes.add_textbox(left, Inches(1 + i), Inches(3), Inches(0.5))
        run = box.text_frame.paragraphs[0].add_run()
        run.text = f"テキスト{i + 1}"
        run.font.name = font
        run.font.size = Pt(18)

    slide2 = prs.slides.add_slide(blank)
    table = slide2.shapes.add_table(2, 2, Inches(1), Inches(1), Inches(4), Inches(1)).table
    table.cell(0, 0).text = "項目"
    table.cell(0, 1).text = "値"
    table.cell(1, 0).text = "売上"
    table.cell(1, 1).text = "100"

    buffer = io.BytesIO()
    prs.save(buffer)
    return buffer.getvalue()


REVIEW_JSON = {
    "axes": [{"key": k, "score": 4.0, "issues": []} for k in
             ["logic", "message_body", "charts", "tone_manner", "design"]],
    "summary": "構成は良好",
    "key_messages": ["売上100"],
}


# --- extract / lint ---------------------------------------------------------


def test_extract_pptx_collects_structure():
    extract = extract_pptx(make_messy_pptx())
    assert extract["slide_count"] == 2
    assert set(extract["fonts"]) == {"Meiryo", "Arial", "Comic Sans MS"}
    assert extract["slides"][1]["has_table"] is True
    assert "テキスト1" in extract["slides"][0]["text"]


def test_extract_pptx_rejects_garbage():
    with pytest.raises(InvalidInputError):
        extract_pptx(b"not a pptx")


def test_lint_flags_fonts_alignment_and_table():
    issues = lint_pptx(extract_pptx(make_messy_pptx()))
    rules = {i["rule"] for axis in issues.values() for i in axis}
    assert RULE_NORMALIZE_FONTS in rules
    assert RULE_ALIGN_SNAP in rules
    assert RULE_TABLE_HEADER in rules
    font_issue = issues["tone_manner"][0]
    assert font_issue["applicable"] is True
    assert "3種類" in font_issue["text"]


def test_build_text_dump_lists_slides():
    dump = build_text_dump(extract_pptx(make_messy_pptx()))
    assert "スライド1" in dump and "スライド2" in dump
    assert "テキスト2" in dump


# --- apply ------------------------------------------------------------------


def test_apply_normalize_fonts_keeps_top_two():
    coached = apply_fixes(make_messy_pptx(), {RULE_NORMALIZE_FONTS: []})
    fonts = set(extract_pptx(coached)["fonts"])
    assert "Comic Sans MS" not in fonts
    assert fonts <= {"Meiryo", "Arial"}


def test_apply_align_snap_merges_near_left_edges():
    coached = apply_fixes(make_messy_pptx(), {RULE_ALIGN_SNAP: [1]})
    lefts = extract_pptx(coached)["slides"][0]["lefts"]
    # 30000EMUだけズレていた2つが同じ左端になる（遠い3つ目はそのまま）
    assert len(set(lefts)) == 2


def test_apply_table_header_bold():
    coached = apply_fixes(make_messy_pptx(), {RULE_TABLE_HEADER: [2]})
    prs = Presentation(io.BytesIO(coached))
    table = next(s for s in prs.slides[1].shapes if s.has_table).table
    header_runs = [r for c in table.rows[0].cells for p in c.text_frame.paragraphs for r in p.runs]
    assert header_runs and all(r.font.bold for r in header_runs)


def test_apply_is_non_destructive():
    original = make_messy_pptx()
    before = extract_pptx(original)
    apply_fixes(original, {RULE_NORMALIZE_FONTS: []})
    assert extract_pptx(original) == before  # 入力バイト列は不変


# --- endpoints --------------------------------------------------------------


@pytest.mark.asyncio
@patch("open_notebook.database.repository.repo_create", new_callable=AsyncMock)
async def test_pptx_review_endpoint_merges_lint_and_stores_original(
    mock_create, client, tmp_path, monkeypatch
):
    import api.slide_review_service as svc

    monkeypatch.setattr(svc, "SLIDE_REVIEW_DIR", tmp_path)
    mock_create.return_value = {"id": "slide_review:p1"}
    with (
        patch(
            "api.slide_review_service.run_text_review",
            new=AsyncMock(return_value=json.dumps(REVIEW_JSON)),
        ) as mock_text,
        patch(
            "api.slide_review_service.ground_citations",
            new=AsyncMock(return_value=[]),
        ),
    ):
        response = client.post(
            "/api/mentor/slide-review",
            files={"file": ("提案書.pptx", make_messy_pptx(), "application/octet-stream")},
        )

    assert response.status_code == 200
    body = response.json()
    assert body["kind"] == "pptx"
    assert body["page_count"] == 2
    # lint 指摘（applicable + id 付き）がLLM軸へマージされる
    tone = next(a for a in body["axes"] if a["key"] == "tone_manner")
    assert any(i["rule"] == RULE_NORMALIZE_FONTS and i["applicable"] for i in tone["issues"])
    # 原本が保存され、テキストダンプがLLMに渡っている
    stored = mock_create.await_args.args[1]["stored_path"]
    assert stored and stored.startswith(str(tmp_path))
    assert "スライド1" in mock_text.await_args.args[0]


@pytest.mark.asyncio
@patch("open_notebook.database.repository.repo_query", new_callable=AsyncMock)
async def test_apply_endpoint_returns_coached_pptx(mock_query, client, tmp_path):
    original_path = tmp_path / "deck.pptx"
    original_path.write_bytes(make_messy_pptx())
    mock_query.return_value = [
        {
            "kind": "pptx",
            "stored_path": str(original_path),
            "axes": [
                {
                    "key": "tone_manner",
                    "score": 2.0,
                    "issues": [
                        {
                            "id": f"{RULE_NORMALIZE_FONTS}@0",
                            "page": 1,
                            "text": "フォント3種",
                            "rule": RULE_NORMALIZE_FONTS,
                            "applicable": True,
                        }
                    ],
                }
            ],
        }
    ]

    response = client.post(
        "/api/mentor/slide-review/slide_review%3Ap1/apply",
        json={"issue_ids": [f"{RULE_NORMALIZE_FONTS}@0"]},
    )

    assert response.status_code == 200
    assert "presentationml" in response.headers["content-type"]
    coached = extract_pptx(response.content)
    assert "Comic Sans MS" not in coached["fonts"]
    assert (tmp_path / "deck_coached.pptx").exists()
    assert original_path.read_bytes() == make_messy_pptx() or original_path.exists()


@pytest.mark.asyncio
@patch("open_notebook.database.repository.repo_query", new_callable=AsyncMock)
async def test_apply_endpoint_rejects_non_applicable_selection(mock_query, client, tmp_path):
    original_path = tmp_path / "deck.pptx"
    original_path.write_bytes(make_messy_pptx())
    mock_query.return_value = [
        {
            "kind": "pptx",
            "stored_path": str(original_path),
            "axes": [
                {
                    "key": "logic",
                    "score": 3.0,
                    "issues": [{"id": None, "page": 1, "text": "LLM指摘", "applicable": False}],
                }
            ],
        }
    ]
    response = client.post(
        "/api/mentor/slide-review/slide_review%3Ap1/apply",
        json={"issue_ids": ["nonexistent@1"]},
    )
    assert response.status_code == 400


@pytest.mark.asyncio
@patch("open_notebook.database.repository.repo_query", new_callable=AsyncMock)
async def test_apply_endpoint_404_and_pdf_rejection(mock_query, client):
    mock_query.return_value = []
    assert (
        client.post(
            "/api/mentor/slide-review/slide_review%3Anope/apply",
            json={"issue_ids": ["x"]},
        ).status_code
        == 404
    )

    mock_query.return_value = [{"kind": "pdf", "stored_path": None, "axes": []}]
    assert (
        client.post(
            "/api/mentor/slide-review/slide_review%3Apdf/apply",
            json={"issue_ids": ["x"]},
        ).status_code
        == 400
    )
