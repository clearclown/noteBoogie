"""pptx structural lint + non-destructive fix application (MENTOR_UI_DESIGN.md §11b).

python-pptx で構造（テキスト・フォント・色・座標・表）を取り出し、
決定的な lint（トンマナ・整列・表）を出す。指摘は rule 付きで、apply で
選択されたものだけを修正した `<name>_coached.pptx` を新規生成する
（元ファイルは変更しない）。レイアウト崩れリスクのある修正（要素の
移動・削除）は対象にしない。
"""

import io
from collections import Counter
from typing import Any, Dict, List

from open_notebook.exceptions import InvalidInputError

# lint しきい値
MAX_FONTS = 2
MAX_COLORS = 4
# 整列スナップ: この範囲内の左端は「揃えたいのに微妙にズレている」とみなす（EMU、0.1インチ）
ALIGN_TOLERANCE_EMU = 91440

RULE_NORMALIZE_FONTS = "normalize_fonts"
RULE_ALIGN_SNAP = "align_snap"
RULE_TABLE_HEADER = "table_header_bold"
APPLICABLE_RULES = (RULE_NORMALIZE_FONTS, RULE_ALIGN_SNAP, RULE_TABLE_HEADER)


def _iter_runs(shape):
    if not getattr(shape, "has_text_frame", False):
        return
    for paragraph in shape.text_frame.paragraphs:
        yield from paragraph.runs


def extract_pptx(data: bytes) -> Dict[str, Any]:
    """スライドごとのテキストと、フォント/色/座標/表の構造サマリを返す。"""
    from pptx import Presentation

    try:
        prs = Presentation(io.BytesIO(data))
    except Exception as e:
        raise InvalidInputError(f"Could not open pptx: {e}") from e

    fonts: Counter = Counter()
    colors: Counter = Counter()
    slides: List[Dict[str, Any]] = []
    for index, slide in enumerate(prs.slides):
        texts: List[str] = []
        lefts: List[int] = []
        has_table = False
        for shape in slide.shapes:
            if getattr(shape, "has_table", False):
                has_table = True
            if shape.left is not None and getattr(shape, "has_text_frame", False):
                if (shape.text_frame.text or "").strip():
                    lefts.append(int(shape.left))
            for run in _iter_runs(shape):
                if run.text.strip():
                    texts.append(run.text)
                if run.font.name:
                    fonts[run.font.name] += 1
                color = run.font.color
                if color and color.type is not None and getattr(color, "rgb", None):
                    colors[str(color.rgb)] += 1
        title = ""
        if slide.shapes.title is not None:
            title = slide.shapes.title.text or ""
        slides.append(
            {
                "index": index,
                "title": title,
                "text": "\n".join(texts),
                "lefts": lefts,
                "has_table": has_table,
            }
        )
    return {
        "slide_count": len(slides),
        "slides": slides,
        "fonts": dict(fonts),
        "colors": dict(colors),
    }


def _misaligned_pages(slides: List[Dict[str, Any]]) -> List[int]:
    """左端が「ほぼ同じだが微妙にズレている」シェイプ群を持つページ。"""
    pages = []
    for slide in slides:
        lefts = sorted(slide["lefts"])
        for a, b in zip(lefts, lefts[1:]):
            if 0 < b - a <= ALIGN_TOLERANCE_EMU:
                pages.append(slide["index"] + 1)
                break
    return pages


def lint_pptx(extract: Dict[str, Any]) -> Dict[str, List[Dict[str, Any]]]:
    """決定的 lint。軸キー → SlideIssue 互換 dict のリスト（rule/applicable 付き）。"""
    issues: Dict[str, List[Dict[str, Any]]] = {"tone_manner": [], "design": [], "charts": []}

    font_names = list(extract["fonts"].keys())
    if len(font_names) > MAX_FONTS:
        main = ", ".join(
            name for name, _ in Counter(extract["fonts"]).most_common(MAX_FONTS)
        )
        issues["tone_manner"].append(
            {
                "id": f"{RULE_NORMALIZE_FONTS}@0",
                "page": 1,
                "text": f"フォントが{len(font_names)}種類使われています（{', '.join(font_names)}）。"
                f"主要{MAX_FONTS}種（{main}）に統一しましょう。",
                "fix": f"全テキストのフォントを {main} に正規化",
                "rule": RULE_NORMALIZE_FONTS,
                "applicable": True,
            }
        )

    if len(extract["colors"]) > MAX_COLORS:
        issues["tone_manner"].append(
            {
                "id": "color_count@0",
                "page": 1,
                "text": f"文字色が{len(extract['colors'])}色あります。強調1色+本文1色を基本に絞りましょう。",
                "fix": "カラーパレットを本文色+強調色（+補助1色）に整理",
                "rule": "color_count",
                "applicable": False,  # 色の自動丸めはブランド色を壊しうるので指摘に留める
            }
        )

    for page in _misaligned_pages(extract["slides"]):
        issues["design"].append(
            {
                "id": f"{RULE_ALIGN_SNAP}@{page}",
                "page": page,
                "text": "左端が微妙にズレたテキスト要素があります。端を揃えると視線誘導が安定します。",
                "fix": "近接する左端座標をスナップして整列",
                "rule": RULE_ALIGN_SNAP,
                "applicable": True,
            }
        )

    for slide in extract["slides"]:
        if slide["has_table"]:
            issues["charts"].append(
                {
                    "id": f"{RULE_TABLE_HEADER}@{slide['index'] + 1}",
                    "page": slide["index"] + 1,
                    "text": "表のヘッダ行を太字にして、データ行との区別を明確にしましょう。",
                    "fix": "ヘッダ行を太字化",
                    "rule": RULE_TABLE_HEADER,
                    "applicable": True,
                }
            )
    return {key: value for key, value in issues.items() if value}


def build_text_dump(extract: Dict[str, Any]) -> str:
    """pptx のテキストを LLM レビュー用に整形する（vision の代わり）。"""
    parts = []
    for slide in extract["slides"]:
        parts.append(f"--- スライド{slide['index'] + 1}: {slide['title'] or '(タイトルなし)'}")
        parts.append(slide["text"] or "(本文なし)")
    return "\n".join(parts)


def apply_fixes(data: bytes, rules_by_page: Dict[str, List[int]]) -> bytes:
    """選択された rule を適用した pptx バイト列を返す（非破壊: 入力は変更しない）。

    rules_by_page: {rule: [対象ページ番号(1始まり)。空リスト=全ページ]}
    """
    from pptx import Presentation

    prs = Presentation(io.BytesIO(data))

    if RULE_NORMALIZE_FONTS in rules_by_page:
        extract = extract_pptx(data)
        ranked = Counter(extract["fonts"]).most_common(MAX_FONTS)
        keep = {name for name, _ in ranked}
        target = ranked[0][0] if ranked else None
        if target:
            for slide in prs.slides:
                for shape in slide.shapes:
                    for run in _iter_runs(shape):
                        if run.font.name and run.font.name not in keep:
                            run.font.name = target

    if RULE_ALIGN_SNAP in rules_by_page:
        pages = set(rules_by_page[RULE_ALIGN_SNAP])
        for index, slide in enumerate(prs.slides):
            if pages and (index + 1) not in pages:
                continue
            shapes = [
                s
                for s in slide.shapes
                if s.left is not None
                and getattr(s, "has_text_frame", False)
                and (s.text_frame.text or "").strip()
            ]
            shapes.sort(key=lambda s: int(s.left))
            cluster: List[Any] = []
            for shape in shapes + [None]:
                if (
                    shape is not None
                    and cluster
                    and int(shape.left) - int(cluster[-1].left) <= ALIGN_TOLERANCE_EMU
                ):
                    cluster.append(shape)
                    continue
                if len(cluster) > 1:
                    snap_to = min(int(s.left) for s in cluster)
                    for member in cluster:
                        member.left = snap_to
                cluster = [shape] if shape is not None else []

    if RULE_TABLE_HEADER in rules_by_page:
        pages = set(rules_by_page[RULE_TABLE_HEADER])
        for index, slide in enumerate(prs.slides):
            if pages and (index + 1) not in pages:
                continue
            for shape in slide.shapes:
                if not getattr(shape, "has_table", False):
                    continue
                for cell in shape.table.rows[0].cells:
                    for paragraph in cell.text_frame.paragraphs:
                        for run in paragraph.runs:
                            run.font.bold = True

    output = io.BytesIO()
    prs.save(output)
    return output.getvalue()
